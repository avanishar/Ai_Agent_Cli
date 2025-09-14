import os
import re
import random
import google.generativeai as genai
from openpyxl import Workbook
from docx import Document
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation

# -------------------------------
# Setup Gemini
# -------------------------------
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
model = genai.GenerativeModel("gemini-1.5-flash")

# -------------------------------
# Utility
# -------------------------------
def sanitize_filename(text: str) -> str:
    return re.sub(r'[^a-zA-Z0-9_\-]', '_', text)[:50]

# -------------------------------
# File Generators
# -------------------------------
def save_articles(prompt: str, output_folder: str):
    os.makedirs(output_folder, exist_ok=True)
    response = model.generate_content(
        f"Write multiple detailed articles based on this prompt:\n{prompt}\n"
        "Format strictly as:\nARTICLE START\n<Title>\n<Content>\nARTICLE END"
    )
    text = response.text.strip()
    articles = re.findall(r"ARTICLE START\n(.*?)\n(.*?)\nARTICLE END", text, re.DOTALL)
    if not articles:
        articles = [("Generated_Article", text)]
    for i, (title, content) in enumerate(articles, start=1):
        safe_title = sanitize_filename(title.strip()) or f"Article_{i}"
        file_path = os.path.join(output_folder, f"{safe_title}.txt")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(f"{title.strip()}\n{'='*len(title.strip())}\n\n{content.strip()}")
        print(f"✅ Saved {file_path}")

def save_notes(prompt: str, output_folder: str):
    os.makedirs(output_folder, exist_ok=True)
    response = model.generate_content(f"Write structured study notes on: {prompt}")
    file_path = os.path.join(output_folder, "notes.txt")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(response.text.strip())
    print(f"✅ Notes saved at {file_path}")

def save_qna(prompt: str, output_folder: str):
    os.makedirs(output_folder, exist_ok=True)
    response = model.generate_content(
        f"Generate 10 question and answer pairs on: {prompt}. "
        "Format as:\nQ: <question>\nA: <answer>\n"
    )
    file_path = os.path.join(output_folder, "QnA.txt")
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(response.text.strip())
    print(f"✅ Q&A saved at {file_path}")

def save_excel(prompt: str, output_folder: str):
    os.makedirs(output_folder, exist_ok=True)
    match = re.search(r"(\d+)", prompt)
    n = int(match.group(1)) if match else 10

    if "employee" in prompt.lower():
        headers = ["ID", "Name", "Department", "Salary"]
        departments = ["HR", "IT", "Finance", "Marketing"]
        rows = [[i+1, f"Employee_{i+1}", random.choice(departments), random.randint(30000, 100000)] for i in range(n)]
        filename = "employees.xlsx"
    elif "student" in prompt.lower():
        headers = ["Roll No", "Name", "Course", "Marks"]
        courses = ["Math", "CS", "Physics", "Biology"]
        rows = [[i+1, f"Student_{i+1}", random.choice(courses), random.randint(40, 100)] for i in range(n)]
        filename = "students.xlsx"
    else:
        headers = ["ID", "Name", "Value"]
        rows = [[i+1, f"Item_{i+1}", random.randint(1, 500)] for i in range(n)]
        filename = "data.xlsx"

    file_path = os.path.join(output_folder, filename)
    wb = Workbook()
    ws = wb.active
    ws.append(headers)
    for row in rows:
        ws.append(row)
    wb.save(file_path)
    print(f"✅ Excel file created at {file_path}")

def save_doc(prompt: str, output_folder: str):
    os.makedirs(output_folder, exist_ok=True)
    response = model.generate_content(f"Write a structured document on: {prompt}")
    file_path = os.path.join(output_folder, "document.docx")
    doc = Document()
    doc.add_heading(prompt, level=1)
    doc.add_paragraph(response.text.strip())
    doc.save(file_path)
    print(f"✅ Word document saved at {file_path}")

def save_pdf(prompt: str, output_folder: str):
    os.makedirs(output_folder, exist_ok=True)
    response = model.generate_content(f"Write a detailed PDF report on: {prompt}")
    file_path = os.path.join(output_folder, "report.pdf")
    doc = SimpleDocTemplate(file_path)
    styles = getSampleStyleSheet()
    elements = [Paragraph(prompt, styles['Heading1']), Spacer(1, 12)]
    for para in response.text.strip().split("\n"):
        elements.append(Paragraph(para, styles['Normal']))
        elements.append(Spacer(1, 12))
    doc.build(elements)
    print(f"✅ PDF saved at {file_path}")

def save_ppt(prompt: str, output_folder: str):
    os.makedirs(output_folder, exist_ok=True)
    response = model.generate_content(
        f"Write bullet points for 5 slides on: {prompt}. "
        "Strict format:\nSLIDE <n>: <Title>\n- point 1\n- point 2\n..."
    )
    file_path = os.path.join(output_folder, "slides.pptx")
    prs = Presentation()

    slides_text = re.findall(
        r"SLIDE\s*\d+:\s*(.*?)\n([\s\S]*?)(?=SLIDE|\Z)",
        response.text, re.MULTILINE
    )

    if not slides_text:  # fallback if parsing fails
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = prompt
        slide.placeholders[1].text = response.text.strip()
    else:
        for title, content in slides_text:
            slide_layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title.strip()
            body = slide.placeholders[1].text_frame
            first = True
            for line in content.strip().split("\n"):
                if line.strip().startswith("-"):
                    if first:
                        body.text = line.strip()[1:].strip()  # replace default text
                        first = False
                    else:
                        p = body.add_paragraph()
                        p.text = line.strip()[1:].strip()
                        p.level = 0

    prs.save(file_path)
    print(f"✅ PPT saved at {file_path}")

# -------------------------------
# Dispatcher
# -------------------------------
def handle_task(task: str, output_folder: str):
    t = task.lower()
    if "article" in t:
        save_articles(task, output_folder)
    elif "note" in t:
        save_notes(task, output_folder)
    elif "q" in t and "a" in t:
        save_qna(task, output_folder)
    elif "excel" in t or "spreadsheet" in t:
        save_excel(task, output_folder)
    elif "doc" in t or "word" in t:
        save_doc(task, output_folder)
    elif "pdf" in t:
        save_pdf(task, output_folder)
    elif "ppt" in t or "presentation" in t or "slides" in t:
        save_ppt(task, output_folder)
    else:
        print("❌ Sorry, I don’t know how to handle this task yet.")

# -------------------------------
# Agent Loop
# -------------------------------
def ai_agent():
    print("🤖 AI Agent Ready. Supports: Articles, Notes, Q&A, Excel, Word Docs, PDF, PPT.\n")
    while True:
        task = input("📝 Enter a task (or 'exit' to quit): ")
        if task.lower() == "exit":
            break
        output_folder = input("📂 Enter output folder (default: outputs): ").strip() or "outputs"
        handle_task(task, output_folder)

if __name__ == "__main__":
    ai_agent()
