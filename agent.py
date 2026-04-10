import os
import re
import json
import datetime
from pathlib import Path
import requests

from openpyxl import Workbook
from docx import Document
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

OUTPUTS_DIR = Path("outputs")
OUTPUTS_DIR.mkdir(exist_ok=True)

# ─── Helpers ─────────────────────────────

def timestamp():
    return datetime.datetime.now().strftime("%m%d_%H%M%S")

def sanitize(text):
    return "_".join(text.split()[:3])

# 🔥 OpenRouter LLM
def llm(prompt):
    response = requests.post(
        "https://openrouter.ai/api/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {os.getenv('OPENROUTER_API_KEY')}",
            "Content-Type": "application/json"
        },
        json={
            "model": "openai/gpt-3.5-turbo",
            "messages": [{"role": "user", "content": prompt}]
        }
    )

    if response.status_code != 200:
        raise Exception(response.text)

    return response.json()["choices"][0]["message"]["content"]


# ─── GENERATORS ─────────────────────────

def generate_notes(topic):
    content = llm(f"Create structured notes on: {topic}")
    filename = f"Notes_{sanitize(topic)}_{timestamp()}.txt"
    path = OUTPUTS_DIR / filename
    path.write_text(content)
    return str(path), filename


def generate_word(topic):
    content = llm(f"Write document on {topic}")

    doc = Document()
    doc.add_heading(topic, 1)

    for line in content.split("\n"):
        doc.add_paragraph(line)

    filename = f"Doc_{sanitize(topic)}_{timestamp()}.docx"
    path = OUTPUTS_DIR / filename
    doc.save(path)

    return str(path), filename


def generate_pdf(topic):
    content = llm(f"Write report on {topic}")

    filename = f"Report_{sanitize(topic)}_{timestamp()}.pdf"
    path = OUTPUTS_DIR / filename

    doc = SimpleDocTemplate(str(path), pagesize=letter)
    styles = getSampleStyleSheet()

    elements = [Paragraph(topic, styles["Heading1"])]

    for line in content.split("\n"):
        elements.append(Paragraph(line, styles["Normal"]))

    doc.build(elements)
    return str(path), filename


def generate_excel(topic):
    wb = Workbook()
    ws = wb.active

    ws.append(["Topic", "Value"])
    for i in range(1, 11):
        ws.append([topic, i])

    filename = f"Data_{sanitize(topic)}_{timestamp()}.xlsx"
    path = OUTPUTS_DIR / filename
    wb.save(path)

    return str(path), filename


def generate_ppt(topic):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Overview"
    slide.placeholders[1].text = llm(f"Give bullet points on {topic}")

    filename = f"PPT_{sanitize(topic)}_{timestamp()}.pptx"
    path = OUTPUTS_DIR / filename
    prs.save(path)

    return str(path), filename


# ─── ROUTER ────────────────────────────

FILE_MAP = {
    "notes": generate_notes,
    "word": generate_word,
    "pdf": generate_pdf,
    "excel": generate_excel,
    "ppt": generate_ppt,
}

def generate_file(file_type, topic):
    return FILE_MAP.get(file_type, generate_notes)(topic)


# ─── CHAT ─────────────────────────────

def chat(message):
    return llm(message)

def chat_stream(message):
    yield llm(message)

def clear_chat():
    pass